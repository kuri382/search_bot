import os
import shutil
import subprocess
import time
from glob import glob

import pptx
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from selenium import webdriver

headers = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_13_6) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/69.0.3497.100 Safari/537.36",
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,image/apng,*/*;q=0.8",
}

def get_html(url, params=None):
    """ get_html
    url: データを取得するサイトのURL
    [params]: 検索サイトのパラメーター {x: param}
    """
    try:
        # データ取得
        resp = requests.get(url, params=params, headers=headers)
        # 要素の抽出
        soup = BeautifulSoup(resp.text, "html.parser")
        return soup
    except Exception as e:
        return None
    
def get_search_url(word, engine="google"):
    """get_search_url
    word: search word
    engine: default google
    """

def make_ppt_file(file_name, slide_titles, slide_urls):
    ppt = Presentation()
    width = ppt.slide_width
    height = ppt.slide_height

    #使用するスライドの種類
    title_slide_layout = ppt.slide_layouts[0] #Title Slideの作成
    bullet_slide_layout = ppt.slide_layouts[1] #Title and Contentの作成
    blank_slide_layout = ppt.slide_layouts[6] #Blankの作成

    #Title Slide
    slide = ppt.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "画像自動まとめ Powerpoint"
    subtitle.text = "python"

    #画像をBlankに張り付ける準備
    fnms = glob('images/*.png')

    tx_left = tx_top = tx_width = tx_height = Inches(0)
    common = Inches(1)

    i = 0
    print(str(len(slide_titles)))

    for fnm in fnms:
        #画像をBlankに張る
        slide_picture = ppt.slides.add_slide(blank_slide_layout)
        #fnm画像，x方向位置,y方向位置，幅，高さ
        pic = slide_picture.shapes.add_picture(fnm, tx_left, tx_top + Inches(0.6), width, height/1.25)

        #ついでに空きスペースにテキストボックスを挿入する
        txBox = slide_picture.shapes.add_textbox(tx_left, tx_top, width, tx_height)
        tB = txBox.text_frame
        print(slide_titles[i])
        txBox.text_frame.add_paragraph().font.size = Pt(8)
        tB.text = slide_titles[i]
        
        txBox_2 = slide_picture.shapes.add_textbox(tx_left, tx_top + Inches(0.3), width, tx_height)
        tB_2 = txBox_2.text_frame
        print(slide_urls[i])
        txBox_2.text_frame.add_paragraph().font.size = Pt(8)
        tB_2.text = slide_urls[i]
        

        i += 1

    ppt.save(file_name)

def screen_shot(chrome_driver, url, number):
    #chromeドライバーにURLの値を入力
    chrome_driver.get(url)
    #chrome windowサイズの決定
    chrome_driver.set_window_size(1250, 1036)
    #chrome 表示サイズズーム
    #driver.execute_script("document.body.style.zoom='90%'")
    # スクリーンショットを指定のタイトルで保存
    chrome_driver.save_screenshot("./images/" + str(number) + ".png")

    return number


def main():
    try:
        #既存のフォルダの削除
        shutil.rmtree('images')
        print('新規フォルダ作成します')
    except Exception as e:
        print("削除するフォルダはありません．新規フォルダ作成します．")
    #新規フォルダの作成
    os.mkdir('images')

    try:
        # urlを代入
        search_url = "https://www.google.co.jp/search"
        search_params = {"q": "コラボ商品　プレスリリース"}
        # データ取得
        time.sleep(1)
        soup = get_html(search_url, search_params)
        # クラスを指定し，aタグのもののみ抽出
        tags = soup.select('.yuRUbf > a')

        #chrome立ち上げ
        driver = webdriver.Chrome()
        #保存する画像ファイル名を初期化
        file_number = 1
        title_lists = []
        url_lists = []

        for tag in tags:
            try:
                #記事タイトルを取得
                current_title= tag.get_text()
                #urlのhref以下を取得
                url = tag.get("href")
                #実行中の記事タイトル，URLを表示
                #print(current_title)
                #print(url)

                screen_shot(driver, url, file_number)

                #ファイル名の繰り上げ
                file_number += 1
                title_lists.append(current_title)
                print(str(len(title_lists)))
                url_lists.append(url)

            except Exception as e:
                print("画像取得ができません")
        
    except Exception as e:
        print("取得できませんでした")

    driver.quit()
    print("making ppt file")
    make_ppt_file("figure1.pptx", title_lists, url_lists)
    print("ppt file has made")

main()
